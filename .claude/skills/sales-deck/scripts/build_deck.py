"""Render a .pptx by duplicating source slides from master.pptx per SLIDE_MAP.

Strategy (v1):
  1. Open assets/master.pptx, which is a clone of the gold-standard sample
     with all 21 source slides intact (original Acefone text preserved).
  2. For each entry in brief.json["slides"], look up its source slide index
     from reference/SLIDE_MAP.md via `TYPE_TO_SOURCE_IDX`.
  3. Duplicate that source slide at the end of the presentation.
  4. Replace the largest-font text frame's text with the brief's slide title.
  5. Write the full brief slide data into the new slide's speaker notes so the
     reviewer has machine-generated content to paste into text frames manually.
  6. After all brief slides are appended, delete the 21 original source slides.
  7. Save to --out.

Usage:
    python scripts/build_deck.py \
        --brief "<path>/brief.json" \
        --spec reference/STYLE_SPEC.md \
        --master assets/master.pptx \
        --out "<output path>/deck.pptx"
"""
from __future__ import annotations

import argparse
import copy
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


# Brief type -> 1-based source-slide index in master.pptx (per SLIDE_MAP.md)
TYPE_TO_SOURCE_IDX = {
    "title": 1,
    "problem_stats": 2,
    "bullet": 3,
    "comparison": 4,
    "positioning": 5,
    "section": 6,
    "trust_badges": 7,
    "stat_grid": 8,
    "feature_grid": 9,
    "use_case_index": 10,
    "use_case_detail": 11,
    "integrations": 16,
    "credibility": 17,
    "results": 18,
    "logos_trusted": 19,
    "quote": 20,
    "closing": 21,
}


WARN_LOG: list[str] = []


def log_warn(msg: str) -> None:
    WARN_LOG.append(msg)
    print(f"WARN: {msg}", file=sys.stderr)


# ---------- I/O ----------


def load_spec(spec_md_path: Path) -> dict:
    text = spec_md_path.read_text(encoding="utf-8")
    blocks = re.findall(r"```json\s*\n(.*?)\n```", text, re.DOTALL)
    if not blocks:
        raise ValueError(f"No ```json``` block found in {spec_md_path}")
    return json.loads(blocks[-1])


def load_brief(brief_path: Path) -> dict:
    return json.loads(brief_path.read_text(encoding="utf-8"))


# ---------- slide duplication ----------


def duplicate_slide(prs, source_slide):
    """Deep-copy a slide's shape tree and relationships into a new slide.

    Returns the new slide. Based on the widely-used python-pptx recipe:
    create a blank slide, remove its empty shape tree, deep-copy elements
    from the source slide's shape tree, then re-register image/media rels.
    """
    blank_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Remove all shapes the blank layout put there
    spTree = new_slide.shapes._spTree  # noqa: SLF001
    for sp in list(spTree):
        tag = sp.tag
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}grpSp") or tag.endswith("}graphicFrame") or tag.endswith("}cxnSp"):
            spTree.remove(sp)

    # Deep-copy each shape element from source
    source_spTree = source_slide.shapes._spTree  # noqa: SLF001
    for el in source_spTree:
        tag = el.tag
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}grpSp") or tag.endswith("}graphicFrame") or tag.endswith("}cxnSp"):
            spTree.append(copy.deepcopy(el))

    # Duplicate part relationships (images, charts, etc.) so the copied XML
    # references remain valid. Skip note-slide rels — notes are slide-scoped.
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part) if not rel.is_external else new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        except Exception:  # noqa: BLE001
            # Fall back to low-level load-rel; may result in duplicate rels but PowerPoint tolerates
            try:
                new_slide.part.rels.add_relationship(rel.reltype, rel.target_part if not rel.is_external else rel.target_ref, rel.rId, rel.is_external)
            except Exception as e:  # noqa: BLE001
                log_warn(f"Could not copy relationship {rel.rId}: {e}")

    return new_slide


def delete_slide_by_id(prs, rId_to_remove: str) -> None:
    """Detach a slide from the presentation by its rId in sldIdLst."""
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    for sldId in list(xml_slides):
        if sldId.get(qn("r:id")) == rId_to_remove:
            xml_slides.remove(sldId)
            # Drop the rel from the presentation part too
            prs.part.drop_rel(rId_to_remove)
            return


def get_slide_rIds_in_order(prs) -> list[str]:
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    return [sldId.get(qn("r:id")) for sldId in xml_slides]


# ---------- text replacement ----------


def replace_largest_text(slide, new_text: str) -> bool:
    """Find the text frame whose largest run has the biggest font size and
    replace its text with new_text. Preserves the existing run's formatting
    by keeping only the first run and dropping the rest.

    Returns True if a replacement happened.
    """
    best_shape = None
    best_size = -1.0
    best_run = None

    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
        except Exception:  # noqa: BLE001
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                size = run.font.size.pt if run.font.size else 0
                if size > best_size:
                    best_size = size
                    best_shape = shape
                    best_run = run

    if best_shape is None or best_run is None or best_size <= 0:
        log_warn(f"No sized text frame found on slide — title '{new_text[:40]}' not written")
        return False

    # Clear the text frame and set the new text in the winning run's paragraph
    tf = best_shape.text_frame
    # Preserve best_run's font attributes
    font_name = best_run.font.name
    font_size = best_run.font.size
    font_bold = best_run.font.bold
    try:
        font_color_rgb = str(best_run.font.color.rgb) if best_run.font.color and best_run.font.color.type is not None else None
    except Exception:  # noqa: BLE001
        font_color_rgb = None

    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = new_text
    if p0.runs:
        r = p0.runs[0]
        if font_name:
            r.font.name = font_name
        if font_size:
            r.font.size = font_size
        if font_bold is not None:
            r.font.bold = font_bold
        if font_color_rgb:
            from pptx.dml.color import RGBColor
            try:
                r.font.color.rgb = RGBColor.from_string(font_color_rgb)
            except Exception:  # noqa: BLE001
                pass

    return True


def attach_brief_to_notes(slide, slide_data: dict) -> None:
    """Dump the full brief slide dict into speaker notes as YAML-ish text."""
    try:
        notes_tf = slide.notes_slide.notes_text_frame
    except Exception as e:  # noqa: BLE001
        log_warn(f"Could not access notes on slide: {e}")
        return

    lines = [f"# brief slide #{slide_data.get('index')} — type: {slide_data.get('type')}"]
    for k, v in slide_data.items():
        if k in {"index", "type"}:
            continue
        if isinstance(v, list):
            lines.append(f"{k}:")
            for item in v:
                lines.append(f"  - {item}")
        else:
            lines.append(f"{k}: {v}")
    notes_tf.text = "\n".join(lines)


# ---------- orchestration ----------


def render(brief: dict, spec: dict, master_path: Path, out_path: Path) -> None:
    if not master_path.exists():
        raise FileNotFoundError(f"master.pptx not found at {master_path} — did Step 0 run?")

    prs = Presentation(str(master_path))
    original_count = len(list(prs.slides))
    original_rIds = get_slide_rIds_in_order(prs)

    source_slides = list(prs.slides)  # snapshot before we append duplicates

    for slide_data in brief.get("slides", []):
        t = slide_data.get("type")
        src_idx = TYPE_TO_SOURCE_IDX.get(t)
        if src_idx is None:
            log_warn(f"Unknown slide type '{t}' at brief index {slide_data.get('index')} — skipped")
            continue
        if src_idx < 1 or src_idx > original_count:
            log_warn(f"Source index {src_idx} for type '{t}' is out of range (master has {original_count}) — skipped")
            continue

        source_slide = source_slides[src_idx - 1]
        new_slide = duplicate_slide(prs, source_slide)

        title = slide_data.get("title") or slide_data.get("heading") or ""
        if title:
            replace_largest_text(new_slide, title)
        attach_brief_to_notes(new_slide, slide_data)

    # Delete the original source slides (they come first in the sldIdLst)
    for rId in original_rIds:
        delete_slide_by_id(prs, rId)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--brief", required=True)
    parser.add_argument("--spec", required=True)
    parser.add_argument("--master", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()

    brief = load_brief(Path(args.brief))
    spec = load_spec(Path(args.spec))
    render(brief, spec, Path(args.master), Path(args.out))

    print("")
    print(f"Rendered: {args.out}")
    print(f"Slides: {len(brief.get('slides', []))}")
    if WARN_LOG:
        print(f"Warnings ({len(WARN_LOG)}):")
        for w in WARN_LOG:
            print(f"  - {w}")
    else:
        print("No warnings.")


if __name__ == "__main__":
    main()
