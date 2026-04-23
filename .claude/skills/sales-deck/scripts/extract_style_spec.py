"""Step 0: extract a style spec + master template from a gold-standard .pptx.

Usage:
    python scripts/extract_style_spec.py \
        --input "<path to sample.pptx>" \
        --out reference/STYLE_SPEC.md \
        --assets assets/

Produces:
    reference/STYLE_SPEC.md      human-readable spec with terminal JSON block
    assets/master.pptx           copy of sample with body-text stripped
    assets/images/               all ppt/media images + manifest.json
    assets/shapes/               sample of recurring non-placeholder shapes
    assets/fonts/                any embedded fonts + README_FONTS.md listing
"""
from __future__ import annotations

import argparse
import hashlib
import io
import json
import shutil
import sys
import zipfile
from collections import Counter, defaultdict
from copy import deepcopy
from datetime import datetime
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ---------- helpers ----------


def emu_to_in(v) -> float:
    if v is None:
        return 0.0
    return round(Emu(v).inches, 2)


def _run_font_info(run) -> dict:
    f = run.font
    color_hex = None
    try:
        if f.color and f.color.type is not None and f.color.rgb is not None:
            color_hex = str(f.color.rgb)
    except Exception:  # noqa: BLE001
        pass
    return {
        "name": f.name,
        "size": float(f.size.pt) if f.size else None,
        "bold": bool(f.bold) if f.bold is not None else None,
        "italic": bool(f.italic) if f.italic is not None else None,
        "color": color_hex,
    }


def _placeholder_role(ph) -> str:
    """Map python-pptx placeholder type to a coarse role."""
    t = ph.placeholder_format.type
    idx = ph.placeholder_format.idx
    # idx == 0 is title, idx == 1 is body/subtitle
    name = str(t).lower() if t is not None else ""
    if "title" in name and "sub" not in name:
        return "title"
    if "subtitle" in name or (idx == 1 and ph.has_text_frame and ph.text_frame.text and len(ph.text_frame.text) < 120):
        return "subtitle"
    if "body" in name or idx == 1:
        return "body"
    if "footer" in name:
        return "footer"
    if "slide number" in name or "slidenumber" in name:
        return "slide_number"
    if "date" in name:
        return "date"
    return name or "other"


def _shape_fill_hex(shape):
    try:
        fill = shape.fill
        if fill.type == 1:  # solid
            return str(fill.fore_color.rgb)
    except Exception:  # noqa: BLE001
        pass
    return None


def _shape_line_hex(shape):
    try:
        line = shape.line
        if line.color and line.color.type is not None and line.color.rgb is not None:
            return str(line.color.rgb)
    except Exception:  # noqa: BLE001
        pass
    return None


# ---------- extraction passes ----------


def extract_theme_colors(pptx_path: Path) -> dict:
    """Parse ppt/theme/theme1.xml and return the color scheme as hex."""
    with zipfile.ZipFile(pptx_path) as z:
        theme_names = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        if not theme_names:
            return {}
        with z.open(sorted(theme_names)[0]) as f:
            tree = etree.parse(f)

    result = {}
    scheme = tree.find(".//a:clrScheme", NS)
    if scheme is None:
        return result
    for child in scheme:
        role = etree.QName(child).localname
        # each role has a single child like a:srgbClr or a:sysClr
        for inner in child:
            tag = etree.QName(inner).localname
            if tag == "srgbClr":
                result[role] = inner.get("val", "").upper()
            elif tag == "sysClr":
                result[role] = inner.get("lastClr", "").upper()
    return result


def extract_embedded_fonts(pptx_path: Path, out_dir: Path) -> list[dict]:
    """Dump embedded font files from ppt/fonts/ and return manifest entries."""
    out_dir.mkdir(parents=True, exist_ok=True)
    fonts = []
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/fonts/") and not name.endswith("/"):
                target = out_dir / Path(name).name
                with z.open(name) as src, open(target, "wb") as dst:
                    shutil.copyfileobj(src, dst)
                fonts.append({"file": target.name, "src_path": name})
    return fonts


def extract_images(pptx_path: Path, out_dir: Path, slide_image_refs: dict) -> list[dict]:
    """Dump ppt/media/* images + write manifest.json."""
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest = []
    with zipfile.ZipFile(pptx_path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/") and not n.endswith("/")]
        for name in media:
            data = z.read(name)
            sha1 = hashlib.sha1(data).hexdigest()[:12]
            basename = Path(name).name
            target = out_dir / basename
            target.write_bytes(data)
            dims = None
            try:
                from PIL import Image  # local import so Pillow is optional at import time
                with Image.open(io.BytesIO(data)) as im:
                    dims = list(im.size)
            except Exception:  # noqa: BLE001
                pass
            manifest.append({
                "file": basename,
                "sha1": sha1,
                "bytes": len(data),
                "dimensions": dims,
                "first_slide": slide_image_refs.get(basename),
            })
    (out_dir / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return manifest


def build_slide_image_refs(pptx_path: Path) -> dict:
    """Map media basename -> first slide index where it appears via rels parsing."""
    refs: dict[str, int] = {}
    with zipfile.ZipFile(pptx_path) as z:
        slide_rels = sorted(
            n for n in z.namelist()
            if n.startswith("ppt/slides/_rels/") and n.endswith(".xml.rels")
        )
        for idx, name in enumerate(slide_rels, start=1):
            with z.open(name) as f:
                tree = etree.parse(f)
            for rel in tree.getroot():
                target = rel.get("Target", "")
                if "media/" in target:
                    base = target.rsplit("/", 1)[-1]
                    refs.setdefault(base, idx)
    return refs


# ---------- main walk ----------


def walk_presentation(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))

    slide_width_in = emu_to_in(prs.slide_width)
    slide_height_in = emu_to_in(prs.slide_height)

    typography = defaultdict(list)   # role -> list[run_font_info]
    color_counter: Counter = Counter()
    line_color_counter: Counter = Counter()
    shape_fingerprints: Counter = Counter()
    slides_info = []

    for i, slide in enumerate(prs.slides):
        s_info = {
            "index": i + 1,
            "layout": slide.slide_layout.name,
            "placeholders": [],
            "shapes": 0,
            "images": 0,
            "text_frames": 0,
            "total_words": 0,
            "word_budget_hint": 0,
        }

        for shape in slide.shapes:
            s_info["shapes"] += 1
            try:
                stype = shape.shape_type
            except (NotImplementedError, Exception):  # noqa: BLE001
                stype = None
            if stype == 13:  # PICTURE
                s_info["images"] += 1

            fill = _shape_fill_hex(shape)
            if fill:
                color_counter[fill.upper()] += 1
            line = _shape_line_hex(shape)
            if line:
                line_color_counter[line.upper()] += 1

            try:
                is_ph = shape.is_placeholder
            except Exception:  # noqa: BLE001
                is_ph = False

            if not is_ph:
                try:
                    w = emu_to_in(shape.width)
                    h = emu_to_in(shape.height)
                    fp = (str(stype), fill, line, w, h)
                    shape_fingerprints[fp] += 1
                except Exception:  # noqa: BLE001
                    pass

            try:
                has_tf = shape.has_text_frame
            except Exception:  # noqa: BLE001
                has_tf = False

            if has_tf:
                s_info["text_frames"] += 1
                role = _placeholder_role(shape) if is_ph else "non_placeholder_text"
                words = 0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        typography[role].append(_run_font_info(run))
                        if run.font.color and run.font.color.type is not None:
                            try:
                                if run.font.color.rgb:
                                    color_counter[str(run.font.color.rgb).upper()] += 1
                            except Exception:  # noqa: BLE001
                                pass
                    words += len((para.text or "").split())
                s_info["total_words"] += words
                if is_ph:
                    s_info["placeholders"].append({
                        "role": role,
                        "idx": shape.placeholder_format.idx,
                        "words": words,
                        "left_in": emu_to_in(shape.left),
                        "top_in": emu_to_in(shape.top),
                        "width_in": emu_to_in(shape.width),
                        "height_in": emu_to_in(shape.height),
                    })

        slides_info.append(s_info)

    # collapse typography: most-common (name, rounded-size, bold) per role
    typo_summary = {}
    for role, runs in typography.items():
        if not runs:
            continue
        key_counter = Counter()
        for r in runs:
            size = round(r["size"]) if r["size"] else None
            key_counter[(r["name"], size, r["bold"])] += 1
        (name, size, bold), count = key_counter.most_common(1)[0]
        typo_summary[role] = {
            "font": name,
            "size_pt": size,
            "bold": bold,
            "sample_count": count,
            "total_runs": len(runs),
        }

    return {
        "slide_width_in": slide_width_in,
        "slide_height_in": slide_height_in,
        "slide_count": len(slides_info),
        "slides": slides_info,
        "typography": typo_summary,
        "shape_fill_colors": dict(color_counter.most_common(40)),
        "shape_line_colors": dict(line_color_counter.most_common(20)),
        "shape_fingerprints": [
            {"shape_type": fp[0], "fill": fp[1], "line": fp[2], "w_in": fp[3], "h_in": fp[4], "count": c}
            for fp, c in shape_fingerprints.most_common(30)
        ],
    }


# ---------- master template ----------


def build_master_template(input_pptx: Path, out_path: Path) -> None:
    """Copy the sample and strip body text from every slide while preserving
    layouts, backgrounds, logos, and shape styling. Placeholders remain so the
    builder can fill them.
    """
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(input_pptx, out_path)
    prs = Presentation(str(out_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                has_tf = shape.has_text_frame
                is_ph = shape.is_placeholder
            except Exception:  # noqa: BLE001
                continue
            # Strip text inside placeholders, keep placeholder structure
            if has_tf and is_ph:
                shape.text_frame.clear()
    prs.save(str(out_path))


# ---------- palette rollup ----------


def build_palette(theme: dict, shape_fills: dict) -> list[dict]:
    palette = []
    role_order = ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                  "dk1", "dk2", "lt1", "lt2"]
    seen = set()
    for role in role_order:
        hx = theme.get(role)
        if hx and hx not in seen:
            palette.append({"role": role, "hex": f"#{hx}", "source": "theme"})
            seen.add(hx)
    for hx, count in shape_fills.items():
        if hx not in seen and len(palette) < 16:
            palette.append({"role": f"shape_fill_{count}x", "hex": f"#{hx}", "source": "shapes"})
            seen.add(hx)
    return palette


# ---------- markdown rendering ----------


def render_markdown(spec: dict, input_pptx: Path) -> str:
    lines = []
    a = lines.append
    a(f"# STYLE_SPEC — extracted from `{input_pptx.name}`")
    a("")
    a(f"_Generated {datetime.now().isoformat(timespec='seconds')} by `extract_style_spec.py`. Re-run to regenerate._")
    a("")
    a(f"Slide size: {spec['slide_width_in']} in × {spec['slide_height_in']} in · total slides: {spec['slide_count']}")
    a("")

    a("## Brand palette")
    a("")
    a("| Role | Hex | Source |")
    a("|---|---|---|")
    for c in spec["palette"]:
        a(f"| {c['role']} | `{c['hex']}` | {c['source']} |")
    a("")

    a("## Typography")
    a("")
    a("| Role | Font | Size (pt) | Bold | Runs sampled |")
    a("|---|---|---|---|---|")
    for role, t in spec["typography"].items():
        a(f"| {role} | {t.get('font') or '—'} | {t.get('size_pt') or '—'} | {t.get('bold')} | {t.get('total_runs')} |")
    a("")

    a("## Slide flow (word-count budgets become render-time limits)")
    a("")
    a("| # | Layout | Placeholders | Words | Shapes | Images |")
    a("|---|---|---|---|---|---|")
    for s in spec["slides"]:
        ph_summary = ", ".join(
            f"{p['role']}({p['words']}w)" for p in s["placeholders"]
        ) or "—"
        a(f"| {s['index']} | {s['layout']} | {ph_summary} | {s['total_words']} | {s['shapes']} | {s['images']} |")
    a("")

    a("## Opening recipe (slide 1)")
    a("")
    if spec["slides"]:
        s0 = spec["slides"][0]
        a(f"- Layout: **{s0['layout']}**")
        a(f"- Shapes: {s0['shapes']}, images: {s0['images']}, total words: {s0['total_words']}")
        for ph in s0["placeholders"]:
            a(f"- Placeholder `{ph['role']}` at ({ph['left_in']}, {ph['top_in']}) in, size {ph['width_in']}×{ph['height_in']} in, {ph['words']} words")
    a("")

    a("## Closing recipe (last slide)")
    a("")
    if spec["slides"]:
        sL = spec["slides"][-1]
        a(f"- Layout: **{sL['layout']}**")
        a(f"- Shapes: {sL['shapes']}, images: {sL['images']}, total words: {sL['total_words']}")
        for ph in sL["placeholders"]:
            a(f"- Placeholder `{ph['role']}` at ({ph['left_in']}, {ph['top_in']}) in, size {ph['width_in']}×{ph['height_in']} in, {ph['words']} words")
    a("")

    a("## Recurring non-placeholder shapes")
    a("")
    a("| Shape type | Fill | Line | Size (in) | Count |")
    a("|---|---|---|---|---|")
    for f in spec["shape_fingerprints"]:
        a(f"| {f['shape_type']} | {f['fill'] or '—'} | {f['line'] or '—'} | {f['w_in']}×{f['h_in']} | {f['count']} |")
    a("")

    a("## Fonts")
    a("")
    if spec["fonts"]:
        for f in spec["fonts"]:
            a(f"- Embedded: `{f['file']}` (from `{f['src_path']}`)")
    else:
        a("- No fonts embedded in source file. Ensure system has the fonts listed in typography table, or add fallbacks below.")
    a("")
    a("### Font fallbacks (edit by hand if system is missing a font)")
    a("")
    a("```json")
    a(json.dumps(spec["font_fallbacks"], indent=2))
    a("```")
    a("")

    a("## Images")
    a("")
    a(f"{len(spec['images'])} images extracted to `assets/images/`. See `assets/images/manifest.json` for full details.")
    a("")

    a("---")
    a("")
    a("## Builder JSON (consumed by `scripts/build_deck.py` — do not edit by hand)")
    a("")
    a("```json")
    builder_spec = {
        "slide_width_in": spec["slide_width_in"],
        "slide_height_in": spec["slide_height_in"],
        "palette": spec["palette"],
        "typography": spec["typography"],
        "font_fallbacks": spec["font_fallbacks"],
        "slide_word_budgets": {
            s["index"]: s["total_words"] for s in spec["slides"]
        },
        "layouts_by_index": {s["index"]: s["layout"] for s in spec["slides"]},
    }
    a(json.dumps(builder_spec, indent=2))
    a("```")
    a("")
    return "\n".join(lines)


# ---------- entry ----------


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True, help="Path to gold-standard .pptx")
    parser.add_argument("--out", required=True, help="Output STYLE_SPEC.md path")
    parser.add_argument("--assets", required=True, help="Output assets/ directory")
    args = parser.parse_args()

    input_pptx = Path(args.input).resolve()
    out_md = Path(args.out).resolve()
    assets_dir = Path(args.assets).resolve()

    if not input_pptx.exists():
        print(f"ERROR: input not found: {input_pptx}", file=sys.stderr)
        sys.exit(2)

    out_md.parent.mkdir(parents=True, exist_ok=True)
    assets_dir.mkdir(parents=True, exist_ok=True)
    (assets_dir / "images").mkdir(exist_ok=True)
    (assets_dir / "shapes").mkdir(exist_ok=True)
    (assets_dir / "fonts").mkdir(exist_ok=True)

    print(f"Parsing presentation: {input_pptx}")
    spec = walk_presentation(input_pptx)

    print("Extracting theme colors…")
    theme = extract_theme_colors(input_pptx)
    spec["palette"] = build_palette(theme, spec["shape_fill_colors"])

    print("Extracting embedded fonts…")
    spec["fonts"] = extract_embedded_fonts(input_pptx, assets_dir / "fonts")

    # Fallback mapping — builder consults this if a font is missing at render time.
    spec["font_fallbacks"] = {}
    for role, info in spec["typography"].items():
        f = info.get("font")
        if f:
            spec["font_fallbacks"].setdefault(f, "Calibri")

    print("Extracting images…")
    slide_image_refs = build_slide_image_refs(input_pptx)
    spec["images"] = extract_images(input_pptx, assets_dir / "images", slide_image_refs)

    print("Building master template (body text stripped)…")
    build_master_template(input_pptx, assets_dir / "master.pptx")

    print("Writing STYLE_SPEC.md…")
    out_md.write_text(render_markdown(spec, input_pptx), encoding="utf-8")

    print("Done.")
    print(f"  STYLE_SPEC: {out_md}")
    print(f"  master.pptx: {assets_dir / 'master.pptx'}")
    print(f"  images: {len(spec['images'])} files in {assets_dir / 'images'}")
    print(f"  fonts: {len(spec['fonts'])} files in {assets_dir / 'fonts'}")
    print(f"  palette entries: {len(spec['palette'])}")


if __name__ == "__main__":
    main()
