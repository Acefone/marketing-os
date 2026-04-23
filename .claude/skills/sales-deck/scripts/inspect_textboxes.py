"""Dump every text frame on specified slides: index, font size of largest run,
position, size, text. Used to calibrate the font-size role heuristic.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def main(pptx_path: str, indices: list[int]) -> None:
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    for idx in indices:
        if idx < 1 or idx > len(slides):
            print(f"[{idx}] OUT OF RANGE")
            continue
        slide = slides[idx - 1]
        print(f"\n=== Slide {idx} (layout={slide.slide_layout.name}) ===")
        rows = []
        for shape_i, shape in enumerate(slide.shapes):
            try:
                if not shape.has_text_frame:
                    continue
            except Exception:  # noqa: BLE001
                continue
            max_pt = 0.0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_pt = max(max_pt, run.font.size.pt)
            text = (shape.text_frame.text or "").strip().replace("\n", " \u21b5 ")
            if len(text) > 80:
                text = text[:80] + "..."
            try:
                left = round(Emu(shape.left).inches, 2) if shape.left else 0
                top = round(Emu(shape.top).inches, 2) if shape.top else 0
                w = round(Emu(shape.width).inches, 2) if shape.width else 0
                h = round(Emu(shape.height).inches, 2) if shape.height else 0
            except Exception:  # noqa: BLE001
                left = top = w = h = 0
            rows.append((shape_i, max_pt, left, top, w, h, text))
        rows.sort(key=lambda r: -r[1])  # largest font first
        for r in rows:
            print(f"  shape[{r[0]:>2}]  {r[1]:>5.1f}pt  @({r[2]:>5},{r[3]:>5})  {r[4]}x{r[5]}  {r[6]}")


if __name__ == "__main__":
    pptx = sys.argv[1]
    ids = [int(x) for x in sys.argv[2:]]
    main(pptx, ids)
