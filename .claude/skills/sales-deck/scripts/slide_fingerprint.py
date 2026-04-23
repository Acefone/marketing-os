"""Print a terse text fingerprint of each slide in a .pptx — to help a human
pick which source-slide index maps to which slide type (title, bullet, stat…).
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


def main(pptx_path: str) -> None:
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        shape_count = 0
        image_count = 0
        for shape in slide.shapes:
            shape_count += 1
            try:
                if shape.shape_type == 13:
                    image_count += 1
            except Exception:  # noqa: BLE001
                pass
            try:
                if shape.has_text_frame:
                    t = (shape.text_frame.text or "").strip()
                    if t:
                        texts.append(t)
            except Exception:  # noqa: BLE001
                continue
        flat = " | ".join(texts)
        if len(flat) > 240:
            flat = flat[:240] + "…"
        print(f"[{i:>2}] shapes={shape_count:>3} images={image_count:>2}  {flat}")


if __name__ == "__main__":
    main(sys.argv[1])
